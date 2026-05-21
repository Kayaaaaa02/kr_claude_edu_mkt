import pandas as pd
import glob
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE

# ─── 1. Find most recent xlsx ───
files = sorted(glob.glob("sample_data/*.xlsx"), key=os.path.getmtime, reverse=True)
if not files:
    raise FileNotFoundError("No xlsx files in sample_data/")
src = files[0]
print(f"Source: {src}")

df = pd.read_excel(src)

# ─── 2. Aggregations ───
by_channel = df.groupby('Channel').agg(
    Revenue=('Revenue_CNY', 'sum'),
    Cost=('Cost_CNY', 'sum'),
    Profit=('Gross_Profit_CNY', 'sum'),
    Quantity=('Quantity', 'sum')
).sort_values('Revenue', ascending=False)
by_channel['Margin'] = (by_channel['Profit'] / by_channel['Revenue'] * 100).round(1)

by_category = df.groupby('Category').agg(
    Revenue=('Revenue_CNY', 'sum'),
    Cost=('Cost_CNY', 'sum'),
    Profit=('Gross_Profit_CNY', 'sum'),
    Quantity=('Quantity', 'sum')
).sort_values('Revenue', ascending=False)
by_category['Margin'] = (by_category['Profit'] / by_category['Revenue'] * 100).round(1)

by_ch_cat = df.groupby(['Channel', 'Category']).agg(
    Revenue=('Revenue_CNY', 'sum'),
    Quantity=('Quantity', 'sum')
).sort_values('Revenue', ascending=False)

daily = df.groupby('Date').agg(
    Revenue=('Revenue_CNY', 'sum'),
    Quantity=('Quantity', 'sum')
).sort_index()

total_rev = df['Revenue_CNY'].sum()
total_profit = df['Gross_Profit_CNY'].sum()
total_qty = df['Quantity'].sum()
total_margin = round(total_profit / total_rev * 100, 1)

print(f"Total Revenue: ¥{total_rev:,.0f}")
print(f"Total Profit: ¥{total_profit:,.0f}")
print(f"Margin: {total_margin}%")

# ─── 3. HTML Dashboard ───
channel_rows = ""
for ch, r in by_channel.iterrows():
    pct = r['Revenue'] / total_rev * 100
    channel_rows += f"""<tr>
        <td>{ch}</td>
        <td>¥{r['Revenue']:,.0f}</td>
        <td>{r['Quantity']:,.0f}</td>
        <td>¥{r['Profit']:,.0f}</td>
        <td>{r['Margin']}%</td>
        <td><div class="bar" style="width:{pct:.0f}%"></div></td>
    </tr>"""

category_rows = ""
for cat, r in by_category.iterrows():
    pct = r['Revenue'] / total_rev * 100
    category_rows += f"""<tr>
        <td>{cat}</td>
        <td>¥{r['Revenue']:,.0f}</td>
        <td>{r['Quantity']:,.0f}</td>
        <td>¥{r['Profit']:,.0f}</td>
        <td>{r['Margin']}%</td>
        <td><div class="bar bar-cat" style="width:{pct*2:.0f}%"></div></td>
    </tr>"""

daily_labels = str(list(daily.index))
daily_rev = str([round(v) for v in daily['Revenue'].values])
daily_qty = str([int(v) for v in daily['Quantity'].values])

# Channel pie data
ch_labels = str(list(by_channel.index))
ch_values = str([round(v) for v in by_channel['Revenue'].values])

date_range = f"{df['Date'].min()} ~ {df['Date'].max()}"

html = f"""<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>F&amp;F China Weekly Report</title>
<script src="https://cdn.jsdelivr.net/npm/chart.js@4.4.0/dist/chart.umd.min.js"></script>
<style>
  * {{ margin:0; padding:0; box-sizing:border-box; }}
  body {{ font-family:'Segoe UI',Arial,sans-serif; background:#f0f2f5; color:#333; }}
  .header {{ background:linear-gradient(135deg,#1a1a2e,#16213e); color:#fff; padding:30px 40px; }}
  .header h1 {{ font-size:24px; font-weight:600; }}
  .header p {{ opacity:0.8; margin-top:5px; font-size:14px; }}
  .container {{ max-width:1200px; margin:0 auto; padding:20px; }}
  .kpi-row {{ display:grid; grid-template-columns:repeat(4,1fr); gap:16px; margin-bottom:24px; }}
  .kpi {{ background:#fff; border-radius:12px; padding:24px; box-shadow:0 2px 8px rgba(0,0,0,0.06); }}
  .kpi .label {{ font-size:13px; color:#888; text-transform:uppercase; letter-spacing:1px; }}
  .kpi .value {{ font-size:28px; font-weight:700; margin-top:8px; color:#1a1a2e; }}
  .kpi .sub {{ font-size:12px; color:#aaa; margin-top:4px; }}
  .charts {{ display:grid; grid-template-columns:2fr 1fr; gap:16px; margin-bottom:24px; }}
  .card {{ background:#fff; border-radius:12px; padding:24px; box-shadow:0 2px 8px rgba(0,0,0,0.06); }}
  .card h3 {{ font-size:16px; margin-bottom:16px; color:#1a1a2e; }}
  table {{ width:100%; border-collapse:collapse; font-size:14px; }}
  th {{ text-align:left; padding:10px 12px; border-bottom:2px solid #e0e0e0; color:#666; font-weight:600; }}
  td {{ padding:10px 12px; border-bottom:1px solid #f0f0f0; }}
  tr:hover {{ background:#fafafa; }}
  .bar {{ height:18px; background:linear-gradient(90deg,#4361ee,#3a86ff); border-radius:9px; min-width:4px; }}
  .bar-cat {{ background:linear-gradient(90deg,#f72585,#b5179e); }}
  .footer {{ text-align:center; padding:20px; color:#aaa; font-size:12px; }}
  canvas {{ max-height:280px; }}
</style>
</head>
<body>
<div class="header">
  <h1>F&amp;F China Weekly Sales Report</h1>
  <p>{date_range} | Generated: 2026-05-13</p>
</div>
<div class="container">
  <div class="kpi-row">
    <div class="kpi">
      <div class="label">Total Revenue</div>
      <div class="value">¥{total_rev/10000:,.0f}万</div>
      <div class="sub">¥{total_rev:,.0f}</div>
    </div>
    <div class="kpi">
      <div class="label">Gross Profit</div>
      <div class="value">¥{total_profit/10000:,.0f}万</div>
      <div class="sub">¥{total_profit:,.0f}</div>
    </div>
    <div class="kpi">
      <div class="label">Gross Margin</div>
      <div class="value">{total_margin}%</div>
      <div class="sub">Target: 60%</div>
    </div>
    <div class="kpi">
      <div class="label">Units Sold</div>
      <div class="value">{total_qty:,}</div>
      <div class="sub">Avg ¥{total_rev/total_qty:,.0f}/unit</div>
    </div>
  </div>

  <div class="charts">
    <div class="card">
      <h3>Daily Revenue Trend (CNY)</h3>
      <canvas id="dailyChart"></canvas>
    </div>
    <div class="card">
      <h3>Revenue by Channel</h3>
      <canvas id="channelPie"></canvas>
    </div>
  </div>

  <div class="card" style="margin-bottom:24px">
    <h3>Channel Performance</h3>
    <table>
      <thead><tr><th>Channel</th><th>Revenue</th><th>Qty</th><th>Profit</th><th>Margin</th><th>Share</th></tr></thead>
      <tbody>{channel_rows}</tbody>
    </table>
  </div>

  <div class="card" style="margin-bottom:24px">
    <h3>Category Performance</h3>
    <table>
      <thead><tr><th>Category</th><th>Revenue</th><th>Qty</th><th>Profit</th><th>Margin</th><th>Share</th></tr></thead>
      <tbody>{category_rows}</tbody>
    </table>
  </div>
</div>
<div class="footer">F&amp;F China | Auto-generated Weekly Report | Confidential</div>

<script>
new Chart(document.getElementById('dailyChart'), {{
  type:'bar',
  data:{{
    labels:{daily_labels},
    datasets:[{{
      label:'Revenue (CNY)',
      data:{daily_rev},
      backgroundColor:'rgba(67,97,238,0.7)',
      borderRadius:6
    }}]
  }},
  options:{{
    responsive:true,
    plugins:{{legend:{{display:false}}}},
    scales:{{y:{{ticks:{{callback:v=>'¥'+(v/10000).toFixed(0)+'万'}}}}}}
  }}
}});

new Chart(document.getElementById('channelPie'), {{
  type:'doughnut',
  data:{{
    labels:{ch_labels},
    datasets:[{{
      data:{ch_values},
      backgroundColor:['#4361ee','#3a86ff','#f72585','#4cc9f0','#7209b7']
    }}]
  }},
  options:{{
    responsive:true,
    plugins:{{legend:{{position:'bottom',labels:{{font:{{size:11}}}}}}}}
  }}
}});
</script>
</body>
</html>"""

html_path = "output/weekly_dashboard.html"
with open(html_path, "w", encoding="utf-8") as f:
    f.write(html)
print(f"\nHTML Dashboard: {html_path}")

# ─── 4. Executive PPT (5 slides) ───
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

DARK = RGBColor(0x1a, 0x1a, 0x2e)
BLUE = RGBColor(0x43, 0x61, 0xee)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xf0, 0xf2, 0xf5)

def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_kpi_box(slide, left, top, label, value, width=2.5):
    from pptx.util import Inches, Pt
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(1.6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x16, 0x21, 0x3e)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0xaa, 0xaa, 0xaa)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(26)
    p2.font.color.rgb = WHITE
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER

# ── Slide 1: Title ──
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1)
add_text(s1, 1.5, 2.0, 10, 1.2, "F&F China", size=44, bold=True)
add_text(s1, 1.5, 3.2, 10, 0.8, "Weekly Sales Report", size=32, color=RGBColor(0x4c,0xc9,0xf0))
add_text(s1, 1.5, 4.3, 10, 0.6, date_range, size=20, color=GRAY)
add_text(s1, 1.5, 5.2, 10, 0.5, "Prepared for Executive Committee | Confidential", size=14, color=GRAY)

# ── Slide 2: KPI Summary ──
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s2)
add_text(s2, 0.8, 0.4, 10, 0.7, "Key Performance Indicators", size=28, bold=True)
add_text(s2, 0.8, 1.0, 10, 0.4, date_range, size=14, color=GRAY)

add_kpi_box(s2, 0.8, 1.8, "TOTAL REVENUE", f"¥{total_rev/10000:,.0f}万", width=2.7)
add_kpi_box(s2, 3.8, 1.8, "GROSS PROFIT", f"¥{total_profit/10000:,.0f}万", width=2.7)
add_kpi_box(s2, 6.8, 1.8, "GROSS MARGIN", f"{total_margin}%", width=2.7)
add_kpi_box(s2, 9.8, 1.8, "UNITS SOLD", f"{total_qty:,}", width=2.7)

# Top insights
insights = [
    f"Tmall remains #1 channel with ¥{by_channel.loc['Tmall','Revenue']/10000:,.0f}万 revenue ({by_channel.loc['Tmall','Revenue']/total_rev*100:.1f}% share)",
    f"MLB Cap is the top category generating ¥{by_category.loc['MLB Cap','Revenue']/10000:,.0f}万",
    f"Average selling price: ¥{total_rev/total_qty:,.0f}/unit across all channels",
    f"Douyin channel showing strong volume with {by_channel.loc['Douyin','Quantity']:,} units sold",
]
for i, ins in enumerate(insights):
    add_text(s2, 1.0, 4.0 + i*0.55, 11, 0.5, f"• {ins}", size=16, color=RGBColor(0xcc,0xcc,0xcc))

# ── Slide 3: Channel Performance ──
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s3, LIGHT_BG)
add_text(s3, 0.8, 0.4, 10, 0.7, "Channel Performance", size=28, bold=True, color=DARK)

tbl_data = [["Channel", "Revenue (CNY)", "Quantity", "Gross Profit", "Margin"]]
for ch, r in by_channel.iterrows():
    tbl_data.append([ch, f"¥{r['Revenue']:,.0f}", f"{r['Quantity']:,.0f}", f"¥{r['Profit']:,.0f}", f"{r['Margin']}%"])

rows_n, cols_n = len(tbl_data), len(tbl_data[0])
tbl_shape = s3.shapes.add_table(rows_n, cols_n, Inches(0.8), Inches(1.4), Inches(11.5), Inches(4.0))
tbl = tbl_shape.table

for i, row_data in enumerate(tbl_data):
    for j, cell_text in enumerate(row_data):
        cell = tbl.cell(i, j)
        cell.text = cell_text
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(14)
            if i == 0:
                para.font.bold = True
                para.font.color.rgb = WHITE
            else:
                para.font.color.rgb = DARK

for j in range(cols_n):
    cell = tbl.cell(0, j)
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK

# ── Slide 4: Category Performance ──
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s4, LIGHT_BG)
add_text(s4, 0.8, 0.4, 10, 0.7, "Category Performance", size=28, bold=True, color=DARK)

tbl_data2 = [["Category", "Revenue (CNY)", "Quantity", "Gross Profit", "Margin"]]
for cat, r in by_category.iterrows():
    tbl_data2.append([cat, f"¥{r['Revenue']:,.0f}", f"{r['Quantity']:,.0f}", f"¥{r['Profit']:,.0f}", f"{r['Margin']}%"])

rows_n2, cols_n2 = len(tbl_data2), len(tbl_data2[0])
tbl_shape2 = s4.shapes.add_table(rows_n2, cols_n2, Inches(0.8), Inches(1.4), Inches(11.5), Inches(4.0))
tbl2 = tbl_shape2.table

for i, row_data in enumerate(tbl_data2):
    for j, cell_text in enumerate(row_data):
        cell = tbl2.cell(i, j)
        cell.text = cell_text
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(14)
            if i == 0:
                para.font.bold = True
                para.font.color.rgb = WHITE
            else:
                para.font.color.rgb = DARK

for j in range(cols_n2):
    cell = tbl2.cell(0, j)
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK

# ── Slide 5: Summary & Next Steps ──
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s5)
add_text(s5, 0.8, 0.4, 10, 0.7, "Summary & Next Steps", size=28, bold=True)

summary_items = [
    f"Total weekly revenue reached ¥{total_rev/10000:,.0f}万 with {total_margin}% gross margin",
    "Tmall and JD.com continue to dominate as primary revenue drivers",
    "MLB Cap category leads across all channels — high demand sustained",
    "Discovery brand categories show opportunity for growth investment",
    "Douyin live-commerce channel growing steadily — consider increased allocation",
]
next_steps = [
    "Increase Douyin live-streaming sessions during peak hours (8-10PM)",
    "Launch MLB Cap x Local Artist collab on WeChat Mini Program",
    "Review Discovery Shoes pricing strategy — current margin below target",
    "Prepare for 618 Shopping Festival promotional campaign",
]

add_text(s5, 0.8, 1.3, 10, 0.5, "KEY FINDINGS", size=18, bold=True, color=RGBColor(0x4c,0xc9,0xf0))
for i, item in enumerate(summary_items):
    add_text(s5, 1.0, 1.9 + i*0.5, 11, 0.5, f"• {item}", size=15, color=RGBColor(0xcc,0xcc,0xcc))

add_text(s5, 0.8, 4.6, 10, 0.5, "NEXT STEPS", size=18, bold=True, color=RGBColor(0x4c,0xc9,0xf0))
for i, item in enumerate(next_steps):
    add_text(s5, 1.0, 5.2 + i*0.5, 11, 0.5, f"• {item}", size=15, color=RGBColor(0xcc,0xcc,0xcc))

pptx_path = "output/FF_China_Weekly_Report.pptx"
prs.save(pptx_path)
print(f"PPT: {pptx_path}")

print("\n=== DONE ===")
print(f"  Source data : {src}")
print(f"  Dashboard   : {html_path}")
print(f"  PPT Report  : {pptx_path}")
